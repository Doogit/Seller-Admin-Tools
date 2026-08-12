import re

import pandas as pd
import pytest
from pptx import Presentation

from core import deck, forecast, store, styles
from core.formatting import fmt_money, parse_money

from test_forecast import make_snapshot, opp


@pytest.fixture
def built(db_path, sample_snapshot):
    buf = deck.build_pptx(sample_snapshot, None, {"team": "Energy", "period": "wk32"},
                          db_path=db_path)
    return Presentation(buf)


def _slide_texts(slide):
    out = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            out.append(shape.text_frame.text)
        elif shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    out.append(cell.text)
    return out


def test_consistency_guard_commit_figure(db_path, sample_snapshot, built):
    """The deck may never disagree with the narrative: parse the commit figure
    back out of slide 2 and compare VALUES (via the shared format->parse pair),
    not formatted strings."""
    rollup = forecast.bucket_rollup(sample_snapshot, db_path=db_path)
    slide2 = built.slides[1]
    commit_cell = None
    for shape in slide2.shapes:
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    if cell.text.startswith("Commit"):
                        commit_cell = cell.text
    assert commit_cell, "commit cell not found by label"
    m = re.search(r"-?\$[\d.,]+[MK]?", commit_cell)
    assert m, commit_cell
    assert parse_money(m.group()) == parse_money(fmt_money(rollup["commit"]))


def test_five_slides_chart_and_draft_footer(built):
    assert len(built.slides) == 5
    assert any(getattr(s, "has_chart", False) for s in built.slides[2].shapes), \
        "slide 3 must contain a native chart"
    for i, slide in enumerate(built.slides):
        assert any(styles.DRAFT_FOOTER in t for t in _slide_texts(slide)), \
            f"slide {i + 1} missing DRAFT footer"


def test_top_deals_table_capped_and_truncated(db_path):
    rows = [
        opp(opportunity_id=f"T{i}", opportunity_name="X" * 60 + str(i),
            amount=1000000.0 + i)
        for i in range(15)
    ]
    sid = make_snapshot(db_path, rows, "w1", "2026-08-11")
    prs = Presentation(deck.build_pptx(sid, None, {}, db_path=db_path))
    for shape in prs.slides[3].shapes:
        if shape.has_table:
            table = shape.table
            assert len(table.rows) <= 11  # header + max 10
            name = table.cell(1, 0).text
            assert len(name) <= styles.OPP_NAME_TRUNCATE
            assert name.endswith("…")


def test_build_md_mirrors_deck(db_path, sample_snapshot):
    rollup = forecast.bucket_rollup(sample_snapshot, db_path=db_path)
    md = deck.build_md(sample_snapshot, None, {"team": "Energy", "period": "wk32"},
                       db_path=db_path)
    assert fmt_money(rollup["commit"]) in md
    for section in ("## Scorecard", "## Pipeline by stage", "## Top deals",
                    "## Risks & asks"):
        assert section in md
    assert styles.DRAFT_FOOTER in md


def test_deck_without_sub_vertical_no_traceback(db_path):
    sid = make_snapshot(db_path, [opp(opportunity_id="A1")], "w1", "2026-08-11")
    assert forecast.sub_vertical_split(sid, db_path=db_path) is None
    md = deck.build_md(sid, None, {}, db_path=db_path)
    assert "Sub-vertical" not in md
    prs = Presentation(deck.build_pptx(sid, None, {}, db_path=db_path))
    assert len(prs.slides) == 5


# --- session-3 forecast additions ---

def test_stage_distribution_with_deltas(db_path):
    prior = make_snapshot(db_path, [
        opp(opportunity_id="A", stage_bucket="early", amount=100.0),
        opp(opportunity_id="B", stage_bucket="mid", amount=200.0),
    ], "w1", "2026-08-01")
    cur = make_snapshot(db_path, [
        opp(opportunity_id="A", stage_bucket="mid", amount=100.0),
        opp(opportunity_id="B", stage_bucket="mid", amount=200.0),
    ], "w2", "2026-08-08")
    dist = forecast.stage_distribution(cur, prior, db_path=db_path).set_index("bucket")
    assert dist.loc["mid", "count"] == 2
    assert dist.loc["mid", "delta_count"] == 1
    assert dist.loc["early", "delta_amount"] == -100.0


def test_owner_rollup_unifies_alias_variants(db_path, sample_path, sample_mapping, stage_map):
    from core import importer
    # "K. Dugas" and "Kevin Dugas" must roll up as one owner row via the
    # alias-normalized owner column written at import time.
    alias_index = {"k dugas": "kevin dugas"}
    raw = sample_path.read_text(encoding="utf-8").replace("Kevin Dugas", "K. Dugas", 3)
    result = importer.import_snapshot(
        raw.encode("utf-8"), sample_mapping, "mdy", stage_map, "w1",
        as_of_date="2026-08-11", db_path=db_path, alias_index=alias_index,
    )
    rollup = forecast.owner_rollup(result.snapshot_id, db_path=db_path)
    owners = rollup["owner"].tolist()
    assert "kevin dugas" in owners
    assert "k dugas" not in owners
    assert len([o for o in owners if "dugas" in o]) == 1


def test_owner_rollup_matches_bucket_rollup_totals(db_path, sample_snapshot):
    per_owner = forecast.owner_rollup(sample_snapshot, db_path=db_path)
    total = forecast.bucket_rollup(sample_snapshot, db_path=db_path)
    assert per_owner["commit"].sum() == pytest.approx(total["commit"])
    assert per_owner["upside"].sum() == pytest.approx(total["upside"])


def test_sub_vertical_split_on_sample(db_path, sample_snapshot):
    sv = forecast.sub_vertical_split(sample_snapshot, db_path=db_path)
    assert sv is not None
    assert set(sv["sub_vertical"]) == {"Power & Utilities", "Oil & Gas", "Pipelines"}
    assert (sv["amount"] > 0).all()


def test_top_deals_joins_flags(db_path, sample_snapshot):
    top = forecast.top_deals(sample_snapshot, db_path=db_path)
    assert len(top) == 10
    assert top.iloc[0]["amount"] >= top.iloc[9]["amount"]
    # the $750K blank-sponsor deal is in the top 10 and carries its flag
    gulf = top[top["opportunity_name"] == "Identity Consolidation"]
    assert not gulf.empty and "no_sponsor" in gulf.iloc[0]["flags"]


def test_deck_with_prior_and_quota(db_path):
    """Exercise the prior-snapshot branches: WoW arrows, the two-series stage
    chart, and the coverage ratio — all dead in the suite without a prior."""
    prior = make_snapshot(db_path, [
        opp(opportunity_id="A", stage_bucket="early", amount=100000.0,
            forecast_category="pipeline"),
    ], "w1", "2026-08-01")
    cur = make_snapshot(db_path, [
        opp(opportunity_id="A", stage_bucket="mid", amount=150000.0,
            forecast_category="upside"),
    ], "w2", "2026-08-08")
    meta = {"team": "E", "period": "wk", "quota": 1_000_000.0}
    md = deck.build_md(cur, prior, meta, db_path=db_path)
    assert "no quota" not in md          # coverage rendered as a ratio
    assert "▲" in md                     # upside rose vs last week
    prs = Presentation(deck.build_pptx(cur, prior, meta, db_path=db_path))
    chart = next(shape.chart for shape in prs.slides[2].shapes
                 if getattr(shape, "has_chart", False))
    assert len(list(chart.plots[0].series)) == 2  # this week + last week


def test_deck_nan_amount_renders_blank_not_dollar_nan(db_path):
    """A blank/NaN amount must render empty, never '$nan' (NaN is truthy)."""
    sid = make_snapshot(db_path, [
        opp(opportunity_id="N", opportunity_name="No Amount Deal", amount=None),
        opp(opportunity_id="H", opportunity_name="Has Amount", amount=200000.0),
    ], "w1", "2026-08-11")
    md = deck.build_md(sid, None, {}, db_path=db_path)
    assert "nan" not in md.lower()
    prs = Presentation(deck.build_pptx(sid, None, {}, db_path=db_path))
    texts = [t for slide in prs.slides for t in _slide_texts(slide)]
    assert not any("nan" in t.lower() for t in texts)
