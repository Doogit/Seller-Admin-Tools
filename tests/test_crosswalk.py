import pandas as pd
import pytest
import yaml

from conftest import REPO_ROOT
from core import crosswalk, importer, ingest, mapping, plan, schema, store

FACTS_CSV = REPO_ROOT / "sample_data" / "account_facts_sample.csv"


@pytest.fixture
def facts_mapping():
    headers = list(ingest.load_csv(FACTS_CSV).columns)
    return mapping.suggest_mapping(headers, schema.ACCOUNT_SCHEMA)


@pytest.fixture
def facts(db_path, facts_mapping):
    result = importer.import_account_facts(FACTS_CSV, facts_mapping, db_path=db_path)
    assert not result.blocking
    return store.load_account_facts(db_path=db_path)


def row(facts, name):
    r = facts[facts["account_name"] == name].iloc[0].to_dict()
    return r


def test_account_schema_mapping_resolves_sample_headers(facts_mapping):
    for field in schema.ACCOUNT_SCHEMA.required:
        assert facts_mapping[field], field
    assert facts_mapping["install_base"] == "Install Base"
    assert facts_mapping["regulatory_scope"] == "Regulatory Scope"
    assert facts_mapping["annual_spend"] == "Annual Spend"


def test_account_facts_replace_semantics(db_path, facts_mapping):
    importer.import_account_facts(FACTS_CSV, facts_mapping, db_path=db_path)
    importer.import_account_facts(FACTS_CSV, facts_mapping, db_path=db_path)
    facts = store.load_account_facts(db_path=db_path)
    assert len(facts) == 5  # re-import replaces, never duplicates


def test_account_facts_replace_overwrites_value(db_path, facts_mapping):
    """Row count alone can't tell REPLACE from IGNORE — assert the stored value
    actually changes on re-upsert of an existing account."""
    importer.import_account_facts(FACTS_CSV, facts_mapping, db_path=db_path)
    before = store.load_account_facts(db_path=db_path)
    target = before.iloc[0].to_dict()
    store.upsert_account_facts(
        pd.DataFrame([{**target, "annual_spend": 424242.0}]), db_path=db_path
    )
    after = store.load_account_facts(db_path=db_path)
    assert len(after) == len(before)  # replaced in place, not appended
    row = after[after["account_name"] == target["account_name"]].iloc[0]
    assert row["annual_spend"] == 424242.0


def test_import_account_facts_blocks_on_missing_required(db_path, facts_mapping):
    """A mapping missing a required ACCOUNT_SCHEMA field must block and write
    nothing, not partially import."""
    broken = {**facts_mapping, "sub_vertical": None}
    result = importer.import_account_facts(FACTS_CSV, broken, db_path=db_path)
    assert result.blocking
    assert result.n_accounts == 0
    assert store.load_account_facts(db_path=db_path).empty


def test_landed_partial_gap_classification(facts):
    meridian = row(facts, "meridian energy")
    gaps = crosswalk.gap_table(meridian)
    by_id = gaps.set_index("obligation_id")
    # Sentinel installed -> SIEM obligation landed
    assert by_id.at["CIP-007-6-R4", "status"] == "landed"
    # Entra ID installed -> identity landed
    assert by_id.at["CIP-005-7-R2", "status"] == "landed"
    # CrowdStrike incumbent -> endpoint partial (displace play)
    assert by_id.at["CIP-007-6-R3", "status"] == "partial"
    assert by_id.at["CIP-007-6-R3", "matched_item"] == "CrowdStrike"
    # nothing for OT security -> gap
    assert by_id.at["CIP-005-7-R1", "status"] == "gap"
    assert {"landed", "partial", "gap"} <= set(gaps["status"])


def test_incumbent_splunk_yields_partial_on_siem(facts):
    gulf = row(facts, "gulf stream petroleum")
    gaps = crosswalk.gap_table(gulf)
    siem = gaps[gaps["capability_category"] == "siem"].iloc[0]
    assert siem["status"] == "partial"
    assert siem["matched_item"] == "Splunk"


def test_empty_install_base_all_gap(facts):
    cascade = row(facts, "cascade power light")
    gaps = crosswalk.gap_table(cascade)
    assert not gaps.empty
    assert set(gaps["status"]) == {"gap"}


def test_unknown_known_gap_id_warns_not_crash(facts):
    bluewater = row(facts, "bluewater utilities")
    gaps = crosswalk.gap_table(bluewater)
    assert any("BADID-999" in w for w in gaps.attrs["warnings"])
    # known_gaps forcing: TransCanyon's patch obligation is a gap even though
    # nothing else changes
    tc = row(facts, "transcanyon pipelines")
    tc_gaps = crosswalk.gap_table(tc).set_index("obligation_id")
    assert tc_gaps.at["TSA-SD-Pipeline-2021-02F-III.D", "status"] == "gap"
    assert tc_gaps.at["TSA-SD-Pipeline-2021-02F-III.D", "matched_item"] \
        == "flagged in known_gaps"


def test_scope_without_reference_entries_warns(facts):
    meridian = row(facts, "meridian energy")  # scope includes state_puc
    gaps = crosswalk.gap_table(meridian)
    assert any("state_puc" in w for w in gaps.attrs["warnings"])


def test_whitespace_math_excludes_unresolvable_and_reports(facts, db_path,
                                                           sample_snapshot):
    cascade = row(facts, "cascade power light")
    gaps = crosswalk.gap_table(cascade)  # all gap: NERC_CIP + IEC_62443 scope
    pipeline = store.get_opportunities(sample_snapshot, db_path=db_path)
    cascade_pipeline = pipeline[pipeline["account_name"] == "cascade power light"]
    ws = crosswalk.whitespace_estimate(gaps, cascade_pipeline)
    # Cascade pipeline: Sentinel 480K (siem gap), Defender for OT 525K
    # (ot gap), Entra ID 72K (identity gap) -> all count
    assert ws["whitespace_amount"] == pytest.approx(480000 + 525000 + 72000)
    assert ws["unresolved_products"] == []
    # nonsense product excluded and reported
    heliograph = pipeline[pipeline["account_name"] == "heliograph solar"]
    ws2 = crosswalk.whitespace_estimate(gaps, heliograph)
    assert "GridWidget Pro" in ws2["unresolved_products"]
    assert ws2["whitespace_amount"] == pytest.approx(110000)  # Defender for Cloud only


def test_uncovered_gaps_listed(facts):
    cascade = row(facts, "cascade power light")
    gaps = crosswalk.gap_table(cascade)
    ws = crosswalk.whitespace_estimate(gaps, pd.DataFrame())
    # zero pipeline -> every gap capability is uncovered
    assert set(ws["uncovered"]["capability_category"]) \
        == set(gaps["capability_category"])


def test_resolve_category_matches_whole_words_not_substrings():
    section = {"identity": ["entra"], "siem": ["sentinel", "microsoft sentinel"]}
    # whole-word / whole-phrase matches still resolve
    assert crosswalk.resolve_category("Entra ID", section) == "identity"
    assert crosswalk.resolve_category("Microsoft Sentinel", section) == "siem"
    # 'entra' embedded in 'central' must not over-match
    assert crosswalk.resolve_category("Central Management Console", section) is None
    assert crosswalk.resolve_category("", section) is None


def test_alias_rename_removes_orphan_facts_row(db_path, facts_mapping):
    # Re-keying a facts row onto a pipeline's canonical name must move it, not
    # leave the old-name row orphaned in the table.
    importer.import_account_facts(FACTS_CSV, facts_mapping, db_path=db_path)
    facts = store.load_account_facts(db_path=db_path)
    old, new = "cascade power light", "cascade p and l"
    renamed = facts[facts["account_name"] == old].copy()
    renamed["account_name"] = new
    store.upsert_account_facts(renamed, db_path=db_path)
    store.delete_account_facts(old, db_path=db_path)
    after = store.load_account_facts(db_path=db_path)
    assert old not in after["account_name"].tolist()   # no orphan
    assert new in after["account_name"].tolist()
    assert len(after) == len(facts)                    # moved, not added


def test_append_product_alias_roundtrip(tmp_path):
    p = tmp_path / "product_map.yaml"
    p.write_text("products:\n  siem:\n    - sentinel\nincumbents: {}\n", encoding="utf-8")
    crosswalk.append_product_alias("GridWidget Pro", "siem", path=p)
    data = yaml.safe_load(p.read_text(encoding="utf-8"))
    assert "gridwidget pro" in data["products"]["siem"]
    assert crosswalk.resolve_category("GridWidget Pro", data["products"]) == "siem"


def test_append_alias_roundtrip(tmp_path):
    p = tmp_path / "aliases.yaml"
    p.write_text("accounts: {}\nowners: {}\n", encoding="utf-8")
    ingest.append_alias("Gulfstream Petroleum", "Gulf Stream Petroleum", p)
    idx = ingest.load_alias_index(p)
    assert ingest.normalize_name("Gulf Stream Petroleum", idx) == "gulfstream petroleum"


def test_plan_compose_and_exports(facts, db_path, sample_snapshot):
    meridian = row(facts, "meridian energy")
    gaps = crosswalk.gap_table(meridian)
    pipeline = store.get_opportunities(sample_snapshot, db_path=db_path)
    mer_pipeline = pipeline[pipeline["account_name"] == "meridian energy"]
    ws = crosswalk.whitespace_estimate(gaps, mer_pipeline)
    sections = plan.compose(meridian, gaps, ws, mer_pipeline)
    assert sections["next_actions"]
    assert any("Displacement play" in a for a in sections["next_actions"])

    md = plan.plan_md(sections, "VERIFY BEFORE USE")
    for part in ("## Obligation → capability → gap", "## Whitespace",
                 "## Next actions", "VERIFY BEFORE USE"):
        assert part in md

    from pptx import Presentation
    prs = Presentation(plan.plan_pptx(sections, "VERIFY BEFORE USE"))
    assert len(prs.slides) == 4
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)
    assert any("VERIFY BEFORE USE" in t for t in texts)


def test_all_gap_account_plan_renders(facts):
    cascade = row(facts, "cascade power light")
    gaps = crosswalk.gap_table(cascade)
    ws = crosswalk.whitespace_estimate(gaps, pd.DataFrame())
    sections = plan.compose(cascade, gaps, ws, pd.DataFrame())
    md = plan.plan_md(sections, "d")
    assert "(no open pipeline for this account)" in md
    from pptx import Presentation
    assert len(Presentation(plan.plan_pptx(sections, "d")).slides) == 4


def test_whitespace_unmeasurable_when_no_product():
    import pandas as pd
    from core import crosswalk
    gap_df = pd.DataFrame(
        [{"obligation_id": "X", "capability_category": "siem",
          "product_label": "P", "status": "gap", "matched_item": ""}],
        columns=crosswalk.GAP_COLUMNS,
    )
    pipeline = pd.DataFrame([
        {"opportunity_name": "D1", "stage_bucket": "mid", "amount": 100000.0, "product": ""},
    ])
    ws = crosswalk.whitespace_estimate(gap_df, pipeline)
    assert ws["unmeasurable"] is True
    assert ws["whitespace_amount"] == 0.0


def test_next_action_cites_pipeline_amount():
    import pandas as pd
    from core import crosswalk, plan
    gap_df = pd.DataFrame(
        [{"obligation_id": "X", "framework": "NERC_CIP", "paraphrase": "p",
          "capability_category": "siem", "product_label": "Sentinel",
          "status": "gap", "matched_item": ""}],
        columns=crosswalk.GAP_COLUMNS,
    )
    whitespace = {
        "whitespace_amount": 100000.0,
        "matched": pd.DataFrame([{"opportunity_name": "D1", "product": "sentinel",
                                  "capability_category": "siem", "amount": 100000.0}]),
        "uncovered": pd.DataFrame(columns=["obligation_id", "capability_category",
                                           "product_label"]),
        "unresolved_products": [], "unmeasurable": False,
    }
    sections = plan.compose({"account_name": "acme"}, gap_df, whitespace, None)
    assert any("open pipeline maps to this gap" in a for a in sections["next_actions"])
    assert any("$100K" in a for a in sections["next_actions"])
