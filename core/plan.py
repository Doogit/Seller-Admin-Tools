"""Account plan composer + .md/.pptx export (reuses the deck.py pattern)."""

from __future__ import annotations

from io import BytesIO

import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt

from core import crosswalk, styles
from core.deck import SLIDE_H, SLIDE_W, _add_footer, _add_title
from core.formatting import fmt_money

MCEM_ACTIONS = {
    "uncovered_gap": "Inspire/Design play",
    "gap_with_pipeline": "Design/Empower — pursue open pipeline",
    "partial": "Displacement play at Empower",
    "landed": "Expand at Realize",
}


def compose(account_row, gap_df: pd.DataFrame, whitespace: dict,
            pipeline_df: pd.DataFrame) -> dict:
    """Assemble the MCEM-structured plan sections from precomputed inputs."""
    install = crosswalk._split(account_row.get("install_base"))
    incumbents = crosswalk._split(account_row.get("incumbent_tools"))
    scopes = crosswalk._split(account_row.get("regulatory_scope"))

    open_pipeline = pd.DataFrame()
    if pipeline_df is not None and not pipeline_df.empty:
        open_pipeline = pipeline_df[
            ~pipeline_df["stage_bucket"].fillna("").isin(["closed_won", "closed_lost"])
        ][["opportunity_name", "stage", "amount", "close_date", "owner", "product"]]

    spend = account_row.get("annual_spend")
    summary_bits = [
        f"Sub-vertical: {account_row.get('sub_vertical') or '—'}",
        f"Annual spend: {fmt_money(float(spend)) if spend and pd.notna(spend) else '—'}",
        f"Agreement ends: {account_row.get('agreement_end_date') or '—'}",
        f"Regulatory scope: {', '.join(scopes) or '—'}",
        f"Open pipeline: {len(open_pipeline)} opportunities, "
        f"{fmt_money(float(open_pipeline['amount'].fillna(0).sum()) if not open_pipeline.empty else 0)}",
    ]

    uncovered_cats = set(whitespace["uncovered"]["capability_category"]) \
        if not whitespace["uncovered"].empty else set()
    actions: list[str] = []
    seen: set[tuple] = set()
    for _, g in gap_df.iterrows():
        cat = g["capability_category"]
        if g["status"] == "gap":
            kind = "uncovered_gap" if cat in uncovered_cats else "gap_with_pipeline"
            detail = (f"{MCEM_ACTIONS[kind]}: {cat} ({g['product_label']})"
                      + (" — no pipeline exists yet" if kind == "uncovered_gap" else ""))
        elif g["status"] == "partial":
            detail = (f"{MCEM_ACTIONS['partial']}: displace {g['matched_item']} "
                      f"with {g['product_label']} ({cat})")
        else:
            detail = f"{MCEM_ACTIONS['landed']}: {g['product_label']} landed ({cat})"
        key = (g["status"], cat)
        if key not in seen:
            seen.add(key)
            actions.append(detail)

    return {
        "account_display": account_row.get("account_name_raw")
        or account_row.get("account_name"),
        "summary": summary_bits,
        "footprint": {"install_base": install, "incumbent_tools": incumbents,
                      "exec_contacts": crosswalk._split(account_row.get("exec_contacts"))},
        "gap_table": gap_df,
        "pipeline": open_pipeline,
        "whitespace": whitespace,
        "next_actions": actions,
        "relationship_map": "(fill in — relationship map stays human)",
    }


def plan_md(sections: dict, disclaimer: str) -> str:
    ws = sections["whitespace"]
    lines = [
        f"# Account plan — {sections['account_display']}",
        "",
        "## Account summary",
        *[f"- {b}" for b in sections["summary"]],
        "",
        "## Current footprint",
        f"- Installed: {'; '.join(sections['footprint']['install_base']) or '(none)'}",
        f"- Incumbents: {'; '.join(sections['footprint']['incumbent_tools']) or '(none)'}",
        f"- Exec contacts: {'; '.join(sections['footprint']['exec_contacts']) or '(none)'}",
        "",
        "## Obligation → capability → gap",
        "| Obligation | Capability | Product | Status | Evidence |",
        "|---|---|---|---|---|",
    ]
    for _, g in sections["gap_table"].iterrows():
        lines.append(
            f"| {g['obligation_id']} | {g['capability_category']} | {g['product_label']} "
            f"| {g['status']} | {g['matched_item']} |"
        )
    lines += ["", "## Open pipeline"]
    if sections["pipeline"].empty:
        lines.append("(no open pipeline for this account)")
    else:
        lines += ["| Opportunity | Stage | Amount | Close | Product |", "|---|---|---|---|---|"]
        for _, r in sections["pipeline"].iterrows():
            lines.append(
                f"| {r['opportunity_name']} | {r['stage']} "
                f"| {fmt_money(r['amount']) if pd.notna(r['amount']) else ''} "
                f"| {r['close_date']} | {r['product'] or ''} |"
            )
    lines += [
        "",
        "## Whitespace",
        f"- Open pipeline against gap capabilities: {fmt_money(ws['whitespace_amount'])}",
    ]
    if not ws["uncovered"].empty:
        lines.append("- Uncovered gaps (no play exists yet):")
        for _, u in ws["uncovered"].iterrows():
            lines.append(f"  - {u['obligation_id']}: {u['capability_category']} "
                         f"({u['product_label']})")
    if ws["unresolved_products"]:
        lines.append(
            "- Excluded from estimate (unresolved products): "
            + ", ".join(ws["unresolved_products"])
        )
    lines += [
        "",
        "## MCEM next actions",
        *[f"- {a}" for a in sections["next_actions"]],
        "",
        "## Relationship map",
        sections["relationship_map"],
        "",
        f"*{disclaimer}*",
        f"*{styles.DRAFT_FOOTER}*",
        "",
    ]
    return "\n".join(lines)


def plan_pptx(sections: dict, disclaimer: str) -> BytesIO:
    """Four slides: summary, gap table, pipeline + whitespace, next actions."""
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    def bullets_box(slide, items, left=0.4, top=1.2, width=12.5, height=5.2):
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width),
                                       Inches(height))
        tf = box.text_frame
        tf.word_wrap = True
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(styles.BODY_SIZE_PT)
            p.font.name = styles.FONT_NAME
        return box

    # 1 — summary
    s = prs.slides.add_slide(blank)
    _add_title(s, f"Account plan — {sections['account_display']}", styles.TITLE_SIZE_PT)
    fp = sections["footprint"]
    bullets_box(s, sections["summary"] + [
        "Installed: " + ("; ".join(fp["install_base"]) or "(none)"),
        "Incumbents: " + ("; ".join(fp["incumbent_tools"]) or "(none)"),
    ], top=1.5)

    # 2 — gap table
    s = prs.slides.add_slide(blank)
    _add_title(s, "Obligation → capability → gap")
    gaps = sections["gap_table"]
    headers = ["Obligation", "Capability", "Product", "Status"]
    table = s.shapes.add_table(
        len(gaps) + 1, len(headers), Inches(0.4), Inches(1.1),
        SLIDE_W - Inches(0.8), Inches(0.3) * (len(gaps) + 1),
    ).table
    for j, h in enumerate(headers):
        table.cell(0, j).text = h
    for i, (_, g) in enumerate(gaps.iterrows(), start=1):
        table.cell(i, 0).text = g["obligation_id"]
        table.cell(i, 1).text = g["capability_category"]
        table.cell(i, 2).text = styles.truncate(g["product_label"], 30)
        table.cell(i, 3).text = g["status"]
    for row in table.rows:
        for cell in row.cells:
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(styles.TABLE_SIZE_PT)
                p.font.name = styles.FONT_NAME

    # 3 — pipeline + whitespace
    s = prs.slides.add_slide(blank)
    _add_title(s, "Open pipeline & whitespace")
    ws = sections["whitespace"]
    items = [f"Whitespace vs gap capabilities: {fmt_money(ws['whitespace_amount'])}"]
    for _, r in sections["pipeline"].head(8).iterrows():
        items.append(
            f"{styles.truncate(r['opportunity_name'])} — {r['stage']}, "
            f"{fmt_money(r['amount']) if pd.notna(r['amount']) else '—'}"
        )
    if not ws["uncovered"].empty:
        items.append("Uncovered gaps (no play exists yet): " + ", ".join(
            ws["uncovered"]["capability_category"].unique()
        ))
    if ws["unresolved_products"]:
        items.append("Excluded (unresolved products): "
                     + ", ".join(ws["unresolved_products"]))
    bullets_box(s, items)

    # 4 — next actions
    s = prs.slides.add_slide(blank)
    _add_title(s, "MCEM next actions")
    bullets_box(s, sections["next_actions"] + [sections["relationship_map"]])

    for slide in prs.slides:
        _add_footer(slide)
    # disclaimer on the gap-table slide
    box = prs.slides[1].shapes.add_textbox(
        Inches(0.4), SLIDE_H - Inches(0.75), SLIDE_W - Inches(0.8), Inches(0.3)
    )
    p = box.text_frame.paragraphs[0]
    p.text = disclaimer
    p.font.size = Pt(styles.FOOTER_SIZE_PT)
    p.font.name = styles.FONT_NAME

    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf
