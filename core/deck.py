"""QBR deck builder — .pptx and .md from snapshot analytics.

All numbers come from core.forecast, the same functions page 1 uses, so the
deck can never disagree with the weekly narrative. Styling constants live in
core/styles.py so template polish later is one-file work.
"""

from __future__ import annotations

import datetime as dt
from io import BytesIO

import pandas as pd
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.util import Inches, Pt

from core import forecast, styles
from core.formatting import fmt_money

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def gather(snapshot_id: int, prior_id: int | None, meta: dict, db_path=None) -> dict:
    """Compute every analytic the deck (and the QBR page) needs, once. risk_flags
    is the expensive call — compute it here and thread it into top_deals so the
    page can render + build the .pptx + .md from a single pass."""
    rollup = forecast.bucket_rollup(snapshot_id, db_path=db_path)
    prior_rollup = forecast.bucket_rollup(prior_id, db_path=db_path) if prior_id else None
    flags = forecast.risk_flags(snapshot_id, db_path=db_path)
    at_risk = forecast.at_risk_total(flags)
    quota = meta.get("quota")
    return {
        "rollup": rollup,
        "prior_rollup": prior_rollup,
        "stage_dist": forecast.stage_distribution(snapshot_id, prior_id, db_path=db_path),
        "top": forecast.top_deals(snapshot_id, db_path=db_path, flags=flags),
        "owner_rollup": forecast.owner_rollup(snapshot_id, db_path=db_path, flags=flags),
        "trend": forecast.snapshot_trend(db_path=db_path, through_id=snapshot_id),
        "sub_vertical": forecast.sub_vertical_split(snapshot_id, db_path=db_path),
        "flags": flags,
        "at_risk": at_risk,
        "coverage": (rollup["total_open"] / quota) if quota else None,
    }


def _arrow(current: float, prior: float | None) -> str:
    if prior is None:
        return ""
    if current - prior > 0.005:
        return f" ▲ {fmt_money(current - prior)}"
    if prior - current > 0.005:
        return f" ▼ {fmt_money(prior - current)}"
    return " ▬ flat"


def add_footer(slide) -> None:
    box = slide.shapes.add_textbox(Inches(0.3), SLIDE_H - Inches(0.4),
                                   SLIDE_W - Inches(0.6), Inches(0.3))
    p = box.text_frame.paragraphs[0]
    p.text = styles.DRAFT_FOOTER
    p.font.size = Pt(styles.FOOTER_SIZE_PT)
    p.font.color.rgb = RGBColor(*styles.MUTED_RGB)
    p.font.name = styles.FONT_NAME


def add_title(slide, text: str, size_pt: int = styles.HEADING_SIZE_PT):
    box = slide.shapes.add_textbox(Inches(0.4), Inches(0.25), SLIDE_W - Inches(0.8),
                                   Inches(0.8))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(size_pt)
    p.font.bold = True
    p.font.name = styles.FONT_NAME
    p.font.color.rgb = RGBColor(*styles.ACCENT_RGB)
    return box


def build_pptx(snapshot_id: int, prior_id: int | None = None, meta: dict | None = None,
               db_path=None, data: dict | None = None) -> BytesIO:
    meta = meta or {}
    d = data if data is not None else gather(snapshot_id, prior_id, meta, db_path=db_path)
    rollup, prior_rollup = d["rollup"], d["prior_rollup"]

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    # Slide 1 — title
    s = prs.slides.add_slide(blank)
    add_title(s, f"{meta.get('team', 'Team')} — Business Review",
               styles.TITLE_SIZE_PT)
    sub = s.shapes.add_textbox(Inches(0.4), Inches(1.4), SLIDE_W - Inches(0.8), Inches(1))
    tf = sub.text_frame
    tf.text = meta.get("period", "")
    p2 = tf.add_paragraph()
    p2.text = f"Generated {dt.date.today().isoformat()} from local snapshot data"
    for p in tf.paragraphs:
        p.font.size = Pt(styles.SUBTITLE_SIZE_PT)
        p.font.name = styles.FONT_NAME

    # Slide 2 — scorecard
    s = prs.slides.add_slide(blank)
    add_title(s, "Scorecard")
    cells = [
        ("Commit", fmt_money(rollup["commit"])
         + _arrow(rollup["commit"], prior_rollup["commit"] if prior_rollup else None)),
        ("Upside", fmt_money(rollup["upside"])
         + _arrow(rollup["upside"], prior_rollup["upside"] if prior_rollup else None)),
        ("Coverage", f"{d['coverage']:.1f}x" if d["coverage"] else "— (no quota)"),
        ("At risk", fmt_money(d["at_risk"])),
    ]
    table = s.shapes.add_table(2, 2, Inches(1.5), Inches(1.5), Inches(10), Inches(4)).table
    for i, (label, value) in enumerate(cells):
        cell = table.cell(i // 2, i % 2)
        cell.text = label
        vp = cell.text_frame.add_paragraph()
        vp.text = value
        cell.text_frame.paragraphs[0].font.size = Pt(styles.BODY_SIZE_PT)
        vp.font.size = Pt(styles.HEADING_SIZE_PT)
        vp.font.bold = True
        for p in cell.text_frame.paragraphs:
            p.font.name = styles.FONT_NAME

    # Slide 3 — stage movement (native chart, no image screenshots)
    s = prs.slides.add_slide(blank)
    add_title(s, "Pipeline by stage")
    dist = d["stage_dist"]
    open_dist = dist[~dist["bucket"].isin(["closed_won", "closed_lost"])]
    cd = CategoryChartData()
    cd.categories = list(open_dist["bucket"])
    cd.add_series("This week ($)", tuple(float(v) for v in open_dist["amount"]))
    two_series = "delta_amount" in open_dist.columns
    if two_series:
        prior_amounts = open_dist["amount"] - open_dist["delta_amount"]
        cd.add_series("Last week ($)", tuple(float(v) for v in prior_amounts))
    chart_frame = s.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1), Inches(1.3), Inches(11), Inches(5), cd
    )
    chart = chart_frame.chart
    if two_series:
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False

    # Slide 4 — top deals (capped rows/cols, fixed font, truncated names)
    s = prs.slides.add_slide(blank)
    add_title(s, "Top deals")
    top = d["top"].head(10)
    headers = ["Opportunity", "Account", "Stage", "Amount", "Close"]
    table = s.shapes.add_table(
        len(top) + 1, len(headers), Inches(0.4), Inches(1.2),
        SLIDE_W - Inches(0.8), Inches(0.35) * (len(top) + 1),
    ).table
    for j, h in enumerate(headers):
        table.cell(0, j).text = h
    for i, (_, r) in enumerate(top.iterrows(), start=1):
        table.cell(i, 0).text = styles.truncate(r["opportunity_name"])
        table.cell(i, 1).text = styles.truncate(r["account_name"], 30)
        table.cell(i, 2).text = r["stage"]
        table.cell(i, 3).text = fmt_money(r["amount"]) if pd.notna(r["amount"]) else ""
        table.cell(i, 4).text = r["close_date"]
    for row in table.rows:
        for cell in row.cells:
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(styles.TABLE_SIZE_PT)
                p.font.name = styles.FONT_NAME

    # Slide 5 — risks & asks
    s = prs.slides.add_slide(blank)
    add_title(s, "Risks & asks")
    risk_box = s.shapes.add_textbox(Inches(0.4), Inches(1.2), Inches(7.5), Inches(5))
    tf = risk_box.text_frame
    tf.word_wrap = True
    flags = d["flags"].sort_values("amount", ascending=False)
    if flags.empty:
        tf.text = "No risk flags this period."
    else:
        first = True
        shown = flags.head(8)
        for _, f in shown.iterrows():
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.text = f"• {styles.truncate(f['opportunity_name'])} " \
                     f"({fmt_money(f['amount']) if pd.notna(f['amount']) else '—'}): {f['evidence']}"
            p.font.size = Pt(styles.BODY_SIZE_PT)
            p.font.name = styles.FONT_NAME
        remainder = len(flags) - len(shown)
        if remainder > 0:
            p = tf.add_paragraph()
            p.text = f"…and {remainder} more flagged deal(s) — see the .md appendix."
            p.font.size = Pt(styles.BODY_SIZE_PT)
            p.font.name = styles.FONT_NAME
            p.font.color.rgb = RGBColor(*styles.MUTED_RGB)
    asks = s.shapes.add_textbox(Inches(8.2), Inches(1.2), Inches(4.6), Inches(5))
    tf = asks.text_frame
    tf.text = "Asks"
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.size = Pt(styles.BODY_SIZE_PT)
    ph = tf.add_paragraph()
    ph.text = "(fill in — this section stays human)"
    ph.font.size = Pt(styles.BODY_SIZE_PT)
    ph.font.color.rgb = RGBColor(*styles.MUTED_RGB)

    for slide in prs.slides:
        add_footer(slide)

    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


def build_md(snapshot_id: int, prior_id: int | None = None, meta: dict | None = None,
             db_path=None, data: dict | None = None) -> str:
    meta = meta or {}
    d = data if data is not None else gather(snapshot_id, prior_id, meta, db_path=db_path)
    rollup, prior_rollup = d["rollup"], d["prior_rollup"]
    lines = [
        f"# {meta.get('team', 'Team')} — Business Review ({meta.get('period', '')})",
        "",
        "## Scorecard",
        f"- Commit: {fmt_money(rollup['commit'])}"
        + _arrow(rollup["commit"], prior_rollup["commit"] if prior_rollup else None),
        f"- Upside: {fmt_money(rollup['upside'])}"
        + _arrow(rollup["upside"], prior_rollup["upside"] if prior_rollup else None),
        "- Coverage: " + (f"{d['coverage']:.1f}x" if d["coverage"] else "— (no quota)"),
        f"- At risk: {fmt_money(d['at_risk'])}",
        "",
    ]
    trend = d.get("trend")
    if trend is not None and len(trend) >= 2:
        lines += ["## Trend (commit / upside / at-risk by week)",
                  "| Week | Commit | Upside | At risk |", "|---|---|---|---|"]
        for _, r in trend.iterrows():
            lines.append(
                f"| {r['label']} ({r['as_of_date']}) | {fmt_money(r['commit'])} "
                f"| {fmt_money(r['upside'])} | {fmt_money(r['at_risk'])} |"
            )
        lines.append("")
    lines += ["## Pipeline by stage"]
    for _, r in d["stage_dist"].iterrows():
        delta = f" (Δ {fmt_money(r['delta_amount'])})" if "delta_amount" in r else ""
        lines.append(f"- {r['bucket']}: {int(r['count'])} deals, {fmt_money(r['amount'])}{delta}")
    lines += ["", "## Top deals", "| Opportunity | Account | Stage | Amount | Close |",
              "|---|---|---|---|---|"]
    for _, r in d["top"].head(10).iterrows():
        lines.append(
            f"| {r['opportunity_name']} | {r['account_name']} | {r['stage']} "
            f"| {fmt_money(r['amount']) if pd.notna(r['amount']) else ''} | {r['close_date']} |"
        )
    owners = d.get("owner_rollup")
    if owners is not None and not owners.empty:
        lines += ["", "## By seller",
                  "| Owner | Commit | Upside | Pipeline | Deals | At risk |",
                  "|---|---|---|---|---|---|"]
        for _, r in owners.iterrows():
            lines.append(
                f"| {r['owner']} | {fmt_money(r['commit'])} | {fmt_money(r['upside'])} "
                f"| {fmt_money(r['pipeline'])} | {int(r['deals'])} | {fmt_money(r['at_risk'])} |"
            )
    if d["sub_vertical"] is not None:
        lines += ["", "## Sub-vertical split"]
        for _, r in d["sub_vertical"].iterrows():
            lines.append(f"- {r['sub_vertical']}: {int(r['count'])} deals, {fmt_money(r['amount'])}")
    lines += ["", "## Risks & asks"]
    flags = d["flags"].sort_values("amount", ascending=False)
    if flags.empty:
        lines.append("No risk flags this period.")
    else:
        for _, f in flags.iterrows():
            lines.append(f"- {f['opportunity_name']} "
                         f"({fmt_money(f['amount']) if pd.notna(f['amount']) else '—'}): {f['evidence']}")
    lines += ["- Asks: (fill in — this section stays human)", "",
              f"*{styles.DRAFT_FOOTER}*", ""]
    return "\n".join(lines)
